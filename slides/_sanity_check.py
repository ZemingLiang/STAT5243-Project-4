"""Sanity check for both decks: open with python-pptx, print slide counts and first-slide title.

Run from the slides/ directory:
    python3 _sanity_check.py
"""

from pptx import Presentation


def first_meaningful_text(slide):
    """Return the first non-empty line of text from any shape on a slide."""
    for shape in slide.shapes:
        if not shape.has_text_frame:
            continue
        text = shape.text_frame.text.strip()
        if not text:
            continue
        # The first paragraph is typically the kicker (e.g. "THE QUESTION") on
        # content slides. Walk the paragraphs to find the first one with > 4
        # chars that isn't all-uppercase (i.e. the actual title).
        for para in text.splitlines():
            para = para.strip()
            if len(para) > 4 and not para.isupper():
                return para
        # Fallback: first non-empty paragraph.
        return text.splitlines()[0].strip()
    return ""


def check(path):
    pres = Presentation(path)
    n = len(pres.slides)
    title = first_meaningful_text(pres.slides[0])
    print(f"  {path:30s}  slides={n:3d}  first-title={title!r}")
    return n


if __name__ == "__main__":
    print("python-pptx sanity check")
    print("-" * 70)
    n_oral = check("oral_10min.pptx")
    n_know = check("knowhow_appendix.pptx")
    print("-" * 70)
    print(f"  Total slides across both decks: {n_oral + n_know}")
    assert n_oral == 12, f"expected 12 oral slides, got {n_oral}"
    assert n_know == 13, f"expected 13 knowhow slides (cover + A1..A12), got {n_know}"
    print("  OK")
